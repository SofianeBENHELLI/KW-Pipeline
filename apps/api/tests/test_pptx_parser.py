"""Tests for the PPTX parser (issue #175).

Mirrors the structure of ``test_docx_parser.py`` so the comparison
between the two formats is one-to-one. Fixtures are built in-memory
via ``python-pptx`` itself; nothing binary is committed.
"""

from __future__ import annotations

import io

from fastapi.testclient import TestClient
from pptx import Presentation as build_pptx
from pptx.util import Inches

from app.dependencies import build_services
from app.main import create_app
from app.schemas.extraction import RawExtraction, RawSection
from app.services.document_parser import Parser
from app.services.document_service import DocumentService
from app.services.parsers import PptxParser
from app.services.parsers.pptx import PPTX_CONTENT_TYPE
from app.services.storage_service import InMemoryStorageService

# A "title and content" layout (index 1 in the default master) lets us
# populate both the title placeholder and a body bullet list in tests
# without hand-crafting layouts.
TITLE_AND_CONTENT_LAYOUT_INDEX = 1


def _make_pptx(slides: list[dict] | None = None) -> bytes:
    """Build a tiny PPTX in-memory and return its bytes for parser tests.

    Each ``slides`` entry accepts:
      - ``title``: optional str, populates the slide's title placeholder.
      - ``bullets``: optional list[str], populates the body placeholder.
      - ``table``: optional list[list[str]], adds a Table shape on the slide.
    """
    presentation = build_pptx()
    layout = presentation.slide_layouts[TITLE_AND_CONTENT_LAYOUT_INDEX]
    for spec in slides or []:
        slide = presentation.slides.add_slide(layout)

        title = spec.get("title")
        if title is not None and slide.shapes.title is not None:
            slide.shapes.title.text = title

        bullets = spec.get("bullets") or []
        if bullets:
            # Find the body placeholder (idx 1 on this layout).
            body_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_placeholder = shape
                    break
            if body_placeholder is not None:
                tf = body_placeholder.text_frame
                tf.text = bullets[0]
                for extra in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = extra

        table_rows = spec.get("table")
        if table_rows:
            rows = len(table_rows)
            cols = len(table_rows[0])
            shape = slide.shapes.add_table(
                rows=rows,
                cols=cols,
                left=Inches(0.5),
                top=Inches(2.5),
                width=Inches(6),
                height=Inches(2),
            )
            tbl = shape.table
            for r, row in enumerate(table_rows):
                for c, cell_text in enumerate(row):
                    tbl.cell(r, c).text = cell_text

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _upload(documents: DocumentService, content: bytes, filename: str = "deck.pptx"):
    return documents.upload(filename, PPTX_CONTENT_TYPE, content)


def test_pptx_parser_each_slide_becomes_one_section_in_order() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    content = _make_pptx(
        slides=[
            {
                "title": "Procurement Policy",
                "bullets": ["Annual review", "Dual approval over 100k"],
            },
            {"title": "Supplier Risk", "bullets": ["Late delivery", "Quality incidents"]},
        ],
    )
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    assert extraction.parser_name == "pptx"
    assert extraction.parser_version == "0.1"
    # One section per slide (no tables in this fixture).
    assert [s.heading for s in extraction.sections] == ["Procurement Policy", "Supplier Risk"]
    assert [s.id for s in extraction.sections] == ["slide-0", "slide-1"]
    # Body bullets land in the section text, joined by newline within the
    # placeholder and double-newline between placeholders.
    assert "Annual review" in extraction.sections[0].text
    assert "Dual approval over 100k" in extraction.sections[0].text
    assert "Late delivery" in extraction.sections[1].text
    # parser_metadata carries slide_index for downstream tooling.
    assert extraction.sections[0].parser_metadata == {"slide_index": "0"}
    assert extraction.sections[1].parser_metadata == {"slide_index": "1"}


def test_pptx_parser_emits_one_source_reference_per_section() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    content = _make_pptx(
        slides=[
            {"title": "Slide A", "bullets": ["Body A"]},
            {"title": "Slide B", "bullets": ["Body B"]},
        ],
    )
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    assert len(extraction.source_references) == len(extraction.sections)
    for index, section in enumerate(extraction.sections):
        ref = extraction.source_references[index]
        assert section.source_reference_ids == [ref.id]
        # PPTX exposes pagination — page_number is 1-indexed slide number.
        assert ref.page_number == index + 1
        assert ref.line_start is None
        assert ref.line_end is None


def test_pptx_parser_synthesizes_heading_when_title_missing() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    # First slide has no title, second has one.
    content = _make_pptx(
        slides=[
            {"title": None, "bullets": ["Body without a title"]},
            {"title": "Has Title", "bullets": ["Body two"]},
        ],
    )
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    headings = [s.heading for s in extraction.sections]
    assert headings == ["Slide 1", "Has Title"]


def test_pptx_parser_renders_table_into_section_after_slide_text() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    content = _make_pptx(
        slides=[
            {
                "title": "Risk Matrix",
                "bullets": ["Quarterly review"],
                "table": [["Risk", "Mitigation"], ["Late delivery", "Penalty"]],
            },
        ],
    )
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    # Two sections: the slide's text first, then the table.
    assert [s.id for s in extraction.sections] == ["slide-0", "slide-0-table-0"]
    assert extraction.sections[0].heading == "Risk Matrix"
    assert extraction.sections[1].heading == "Risk Matrix — Table 1"
    assert extraction.sections[1].text == "Risk\tMitigation\nLate delivery\tPenalty"
    assert extraction.sections[1].parser_metadata == {
        "slide_index": "0",
        "table_index": "0",
    }
    # Table reference inherits the slide's page_number.
    assert extraction.source_references[1].page_number == 1


def test_pptx_parser_skips_slides_with_no_text_or_tables() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    # Slide 1 has content; slide 2 has neither title nor body — should
    # produce no section but still be counted as slide 2 (so slide 3's
    # page_number is 3, not 2).
    content = _make_pptx(
        slides=[
            {"title": "Intro", "bullets": ["Welcome"]},
            {"title": None, "bullets": []},
            {"title": "Wrap up", "bullets": ["Thank you"]},
        ],
    )
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    assert [s.heading for s in extraction.sections] == ["Intro", "Wrap up"]
    page_numbers = [r.page_number for r in extraction.source_references]
    assert page_numbers == [1, 3]


def test_pptx_parser_empty_deck_warns_without_raising() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    content = _make_pptx()  # No slides at all.
    version = _upload(documents, content)

    extraction = PptxParser().parse(version=version, storage=documents.storage)

    assert extraction.sections == []
    assert extraction.source_references == []
    assert "No slides with extractable text or tables." in extraction.warnings


def test_pptx_parser_round_trips_through_json() -> None:
    documents = DocumentService(storage=InMemoryStorageService())
    content = _make_pptx(
        slides=[{"title": "Hello", "bullets": ["One", "Two"], "table": [["a", "b"]]}],
    )
    version = _upload(documents, content)
    extraction = PptxParser().parse(version=version, storage=documents.storage)

    payload = extraction.model_dump_json()
    restored = RawExtraction.model_validate_json(payload)

    assert restored == extraction
    assert all(isinstance(section, RawSection) for section in restored.sections)


def test_pptx_parser_conforms_to_parser_protocol() -> None:
    parser = PptxParser()

    assert isinstance(parser, Parser)
    assert parser.supported_content_types == frozenset({PPTX_CONTENT_TYPE})


def test_registry_resolves_pptx_content_type_to_pptx_parser() -> None:
    services = build_services()

    parser = services.parsers.for_content_type(PPTX_CONTENT_TYPE)

    assert isinstance(parser, PptxParser)


def test_pptx_upload_extract_http_flow(monkeypatch) -> None:
    """End-to-end: HTTP upload + extract for the PPTX MIME, with the
    content-type allowlist widened via env var."""
    monkeypatch.setenv("ALLOWED_CONTENT_TYPES", f"text/plain,{PPTX_CONTENT_TYPE}")
    client = TestClient(create_app())

    content = _make_pptx(
        slides=[
            {"title": "Alpha", "bullets": ["Body alpha"]},
            {"title": "Beta", "bullets": ["Body beta"]},
        ],
    )

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("deck.pptx", content, PPTX_CONTENT_TYPE)},
    )
    assert upload_response.status_code == 200, upload_response.text
    version = upload_response.json()
    assert version["content_type"] == PPTX_CONTENT_TYPE

    extract_response = client.post(
        f"/documents/{version['document_id']}/versions/{version['id']}/extract",
    )
    assert extract_response.status_code == 200, extract_response.text
    extraction = extract_response.json()
    assert extraction["parser_name"] == "pptx"
    headings = [section["heading"] for section in extraction["sections"]]
    assert headings == ["Alpha", "Beta"]
    assert len(extraction["source_references"]) == 2
    assert extraction["source_references"][0]["page_number"] == 1
    assert extraction["source_references"][1]["page_number"] == 2
