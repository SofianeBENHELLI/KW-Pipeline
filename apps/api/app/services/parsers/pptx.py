"""PPTX parser backed by python-pptx.

Closes #175. Conforms to the ``Parser`` Protocol declared in
``app.services.document_parser`` so the ``ParserRegistry`` can dispatch
``application/vnd.openxmlformats-officedocument.presentationml.presentation``
uploads to it without any service-layer changes.

python-pptx is BSD-licensed (the no-AGPL/no-GPL constraint stays
satisfied); the dependency is pinned in ``apps/api/pyproject.toml``.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation as load_pptx

from app.schemas.document import DocumentVersion
from app.schemas.extraction import RawExtraction, RawSection, SourceReference
from app.services.storage_service import StorageService

PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _slide_title(slide: Any) -> str | None:
    """Return the slide's title-placeholder text, stripped, or ``None``.

    PowerPoint files where the title placeholder is missing or empty are
    common (custom layouts, stand-alone bullet slides). Callers fall back to
    a synthesized ``"Slide N"`` heading in that case.
    """
    title = slide.shapes.title
    if title is None:
        return None
    text = (title.text or "").strip()
    return text or None


def _shape_text(shape: Any) -> str:
    """Extract a shape's text frame contents, joining paragraphs with ``\\n``.

    Mirrors the DOCX parser's per-paragraph treatment but at the shape level
    — within a single text frame, runs are concatenated and paragraphs are
    line-separated, so a bulleted list keeps its rough visual shape in the
    section's ``text`` field.
    """
    if not shape.has_text_frame:
        return ""
    paragraphs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        line = "".join(run.text for run in paragraph.runs)
        paragraphs.append(line)
    return "\n".join(paragraphs)


def _table_to_text(table: Any) -> str:
    """Render a python-pptx table into the same ``\\t``/``\\n`` shape the
    DOCX and PDF parsers emit, so downstream tooling can stay format-agnostic.
    """
    rows: list[str] = []
    for row in table.rows:
        cells = [(cell.text or "").strip() for cell in row.cells]
        rows.append("\t".join(cells))
    return "\n".join(rows)


class PptxParser:
    """Parser for ``.pptx`` PowerPoint presentations.

    Each slide becomes one text ``RawSection`` (when the slide has any
    extracted text) followed by one ``RawSection`` per ``Table`` shape on
    that slide. Sections are emitted in slide order — slide 1 text, slide 1
    tables, slide 2 text, slide 2 tables, and so on. The slide's title
    placeholder, when present and non-empty, becomes the section heading;
    otherwise the heading falls back to ``"Slide N"`` (1-indexed).

    ``page_number`` on each ``SourceReference`` is 1-indexed and matches the
    slide number a reviewer sees in PowerPoint, giving the review surface a
    "go to slide 7" affordance for free.

    ``parser_metadata`` always carries ``slide_index`` (0-indexed integer
    string), and additionally ``table_index`` for table sections.

    Speaker notes are intentionally not captured in v1 — the DOCX parser
    similarly does not capture footers / headers, and the issue (#175)
    scopes them out. They can be added in a follow-up without changing the
    section schema.
    """

    name = "pptx"
    # Tied to python-pptx's API; bump alongside the dependency floor.
    version = "0.1"
    supported_content_types: frozenset[str] = frozenset({PPTX_CONTENT_TYPE})

    def parse(self, version: DocumentVersion, storage: StorageService) -> RawExtraction:
        """Read stored PPTX bytes and return a typed ``RawExtraction``."""
        content = storage.get(version.storage_uri)
        presentation = load_pptx(io.BytesIO(content))

        sections: list[RawSection] = []
        source_references: list[SourceReference] = []
        text_parts: list[str] = []

        for slide_index, slide in enumerate(presentation.slides):
            page_number = slide_index + 1
            heading = _slide_title(slide) or f"Slide {page_number}"

            # Collect text from every non-table shape in slide-shape order,
            # then emit a single section per slide. Per-shape granularity
            # is intentionally avoided — slides are the natural unit of
            # provenance for a presentation, mirroring how PDF treats
            # pages.
            slide_text_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = _shape_text(shape).strip()
                if text:
                    slide_text_parts.append(text)
            slide_text = "\n\n".join(slide_text_parts)

            if slide_text:
                section_id = f"slide-{slide_index}"
                reference = SourceReference(
                    document_version_id=version.id,
                    section_id=section_id,
                    page_number=page_number,
                    line_start=None,
                    line_end=None,
                    snippet=slide_text[:240],
                )
                section = RawSection(
                    id=section_id,
                    heading=heading,
                    text=slide_text,
                    source_reference_ids=[reference.id],
                    parser_metadata={"slide_index": str(slide_index)},
                )
                sections.append(section)
                source_references.append(reference)
                text_parts.append(slide_text)

            # Tables on this slide, after the slide-text section.
            # ``table_index`` is local to the slide so the section IDs stay
            # short and unique across the whole deck.
            table_index_in_slide = 0
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                table_text = _table_to_text(shape.table)
                section_id = f"slide-{slide_index}-table-{table_index_in_slide}"
                reference = SourceReference(
                    document_version_id=version.id,
                    section_id=section_id,
                    page_number=page_number,
                    line_start=None,
                    line_end=None,
                    snippet=table_text[:240],
                )
                section = RawSection(
                    id=section_id,
                    heading=f"{heading} — Table {table_index_in_slide + 1}",
                    text=table_text,
                    source_reference_ids=[reference.id],
                    parser_metadata={
                        "slide_index": str(slide_index),
                        "table_index": str(table_index_in_slide),
                    },
                )
                sections.append(section)
                source_references.append(reference)
                text_parts.append(table_text)
                table_index_in_slide += 1

        warnings: list[str] = []
        if not sections:
            warnings.append("No slides with extractable text or tables.")

        return RawExtraction(
            document_version_id=version.id,
            parser_name=self.name,
            parser_version=self.version,
            text="\n\n".join(text_parts),
            sections=sections,
            source_references=source_references,
            warnings=warnings,
        )
